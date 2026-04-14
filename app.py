import io
import re
from typing import Dict, List, Tuple

import pandas as pd
import plotly.express as px
import streamlit as st

# NEW: libraries for building PPTX and PDF
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

from reportlab.lib.pagesizes import A4
from reportlab.lib import colors
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
from reportlab.lib.styles import getSampleStyleSheet

# ──────────────────────────────────────────────────────────────────────────────
# App metadata
# ──────────────────────────────────────────────────────────────────────────────
st.set_page_config(page_title="Daily Evaluation Summarizer App (V3)", layout="wide")
st.title("📊 Project DESA — Daily Evaluation Summarizer App (V3)")
st.caption("Project DESA is an application tool used in instantly summarizing the results of the Daily Evaluations during trainings. Developed and enhanced by the Mr. Samson G. Capinig, SEPS of SMME Section of SDO Masbate City.")

# ──────────────────────────────────────────────────────────────────────────────
# Constants: New Template Category Prefixes
# (we match columns that START with these prefixes)
# ──────────────────────────────────────────────────────────────────────────────
CATEGORY_PREFIXES = {
    "PROGRAM MANAGEMENT": "Q06_PROGRAM MANAGEMENT",
    "ACCOMMODATION": "Q07_ACCOMMODATION",
    "TRAINING VENUE": "Q08_TRAINING VENUE",
    "FOOD/MEALS": "Q09_FOOD/MEALS",
    "ADMINISTRATIVE ARRANGEMENTS": "Q10_ADMINISTRATIVE ARRANGEMENTS",
}

# Regex to detect session columns and capture DAY and LM numbers (tolerant to spaces/dashes).
# Handles headers like:
#   Q11_DAY2-LM1->(Session) ...
#   Q12_DAY 1-LM3->(Facilitator) ...
SESSION_REGEX = re.compile(
    r"Q\d+[\s_\-]*DAY\s*(\d+)\s*[-–]?\s*LM\s*(\d+)",
    flags=re.IGNORECASE
)

# ──────────────────────────────────────────────────────────────────────────────
# Helpers
# ──────────────────────────────────────────────────────────────────────────────
def _strip_unnamed(df: pd.DataFrame) -> pd.DataFrame:
    """Drop Unnamed columns often introduced by Excel/CSV exports."""
    return df.loc[:, ~df.columns.astype(str).str.match(r"^Unnamed", na=False)]


def _standardize_headers(df: pd.DataFrame) -> pd.DataFrame:
    """
    Create a slightly cleaned header set for parsing.
    We keep original names for display, but parsing uses a normalized version.
    """
    df = df.copy()
    df.columns = [str(c) for c in df.columns]
    return df


@st.cache_data(show_spinner=False)
def load_file(uploaded_file) -> pd.DataFrame:
    """Read CSV/XLSX and remove unnamed columns."""
    name = uploaded_file.name.lower()
    if name.endswith(".csv"):
        df = pd.read_csv(uploaded_file, dtype=str)  # read raw as string first; coerce later
    else:
        # Read first sheet by default; users can re-upload if they need another sheet
        df = pd.read_excel(uploaded_file, dtype=str)
    df = _strip_unnamed(_standardize_headers(df))
    return df


def coerce_numeric(df: pd.DataFrame) -> pd.DataFrame:
    """Coerce numeric cells to float, non-numeric to NaN."""
    return df.apply(pd.to_numeric, errors="coerce")


def find_category_columns(df: pd.DataFrame) -> Dict[str, List[str]]:
    """
    Map each category to the list of matching columns based on NEW TEMPLATE prefixes.
    """
    cols_by_cat = {cat: [] for cat in CATEGORY_PREFIXES.keys()}
    for col in df.columns:
        col_u = str(col).strip()
        for cat, prefix in CATEGORY_PREFIXES.items():
            if col_u.startswith(prefix):
                cols_by_cat[cat].append(col)
                break
    # Remove empty categories
    return {k: v for k, v in cols_by_cat.items() if v}


def find_session_groups(df: pd.DataFrame) -> Dict[str, List[str]]:
    """
    Group session columns by DAY–LM key, e.g., 'DAY2-LM1'.
    We include both (Session) and (Facilitator) questions under the same group.
    """
    groups: Dict[str, List[str]] = {}
    for col in df.columns:
        match = SESSION_REGEX.search(str(col))
        if match:
            day, lm = match.group(1), match.group(2)
            key = f"DAY{day}-LM{lm}"
            groups.setdefault(key, []).append(col)
    return groups


def compute_avg_for_columns(df: pd.DataFrame, columns: List[str]) -> float:
    """Compute mean across stacked values of the given columns."""
    if not columns:
        return float("nan")
    sub = coerce_numeric(df[columns])
    stacked = sub.stack(dropna=True)  # 1D series of numeric values
    return float(stacked.mean()) if not stacked.empty else float("nan")


def summarize_categories(df: pd.DataFrame, file_label: str) -> pd.DataFrame:
    """
    Return a DataFrame indexed by category, with one column named file_label,
    containing the average for that category.
    """
    colmap = find_category_columns(df)
    stats = {}
    for cat, cols in colmap.items():
        avg = compute_avg_for_columns(df, cols)
        if pd.notna(avg):
            stats[cat] = {file_label: round(avg, 2)}
    return pd.DataFrame(stats).T if stats else pd.DataFrame()


def summarize_sessions(df: pd.DataFrame, file_label: str) -> pd.DataFrame:
    """
    Return a DataFrame indexed by 'DAYx-LMy', with one column named file_label,
    containing the average across all items (Session + Facilitator) for that DAY–LM.
    """
    groups = find_session_groups(df)
    stats = {}
    for session_key, cols in groups.items():
        avg = compute_avg_for_columns(df, cols)
        if pd.notna(avg):
            stats[session_key] = {file_label: round(avg, 2)}
    return pd.DataFrame(stats).T if stats else pd.DataFrame()


def finalize_summary(df: pd.DataFrame) -> pd.DataFrame:
    """Add 'Overall Avg' across all file columns."""
    if df.empty:
        return df
    out = df.copy()
    numeric = out.select_dtypes(include="number")
    overall = numeric.mean(axis=1).round(2)
    out["Overall Avg"] = overall
    return out


def style_numeric(df: pd.DataFrame):
    """Style numeric columns to two decimals."""
    return df.style.format("{:.2f}")


def to_csv_bytes(df: pd.DataFrame) -> bytes:
    return df.to_csv(index=True).encode("utf-8")


def to_excel_bytes(sheets: Dict[str, pd.DataFrame]) -> bytes:
    buf = io.BytesIO()
    with pd.ExcelWriter(buf, engine="openpyxl") as writer:
        for sheet_name, d in sheets.items():
            d.to_excel(writer, sheet_name=sheet_name, index=True)
    buf.seek(0)
    return buf.getvalue()


# ──────────────────────────────────────────────────────────────────────────────
# NEW: Helpers for Reports (PPTX & PDF)
# ──────────────────────────────────────────────────────────────────────────────

DAY_NAME_RE = re.compile(r"DAY[\s_\-]?(\d+)", re.IGNORECASE)  # detect day number in filenames

def _extract_day_from_filename(name: str) -> int | None:
    m = DAY_NAME_RE.search(name or "")
    return int(m.group(1)) if m else None

def _day_headers_from_files(file_names: List[str]) -> List[str]:
    """
    Build ordered day headers like ['Day 1','Day 2',...] from uploaded filenames.
    Fallback to original file names if no day number is found.
    """
    pairs: List[Tuple[str, str, int | None]] = []
    for fn in file_names:
        d = _extract_day_from_filename(fn)
        pairs.append((fn, f"Day {d}" if d else fn, d))
    # Sort by numeric day if available, else keep original order
    have_num = [p for p in pairs if p[2] is not None]
    no_num = [p for p in pairs if p[2] is None]
    have_num.sort(key=lambda x: x[2])
    ordered = have_num + no_num
    return [p[1] for p in ordered]

def _ordered_file_columns(cat_df: pd.DataFrame, file_names: List[str]) -> List[str]:
    """
    Return the columns in the same order as detected days, excluding 'Overall Avg'.
    """
    wanted = []
    # Keep only those file columns that exist in df, in the given order
    for fn in file_names:
        if fn in cat_df.columns:
            wanted.append(fn)
    return wanted

def _map_file_to_day_label(file_names: List[str]) -> Dict[str, str]:
    """
    Map each file column to display header (e.g., 'Day 1').
    """
    mapping = {}
    labels = _day_headers_from_files(file_names)
    for fn, label in zip(file_names, labels):
        mapping[fn] = label
    return mapping

def build_pptx_report(
    category_summary: pd.DataFrame,
    session_summary: pd.DataFrame,
    original_file_names: List[str],
    meta: Dict[str, str],
    highest_score: int = 5,
) -> bytes:
    """
    Create a PPTX presentation following your sample structure:
      - Title
      - Category table with Day columns + Average
      - Session tables per day
      - Overall rating slides
    """
    prs = Presentation()  # blank presentation

    # ---- Title slide (mirrors your sample)
    def add_title_slide():
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        left, top, width, height = Inches(0.5), Inches(0.6), Inches(9), Inches(1.2)
        title = slide.shapes.add_textbox(left, top, width, height).text_frame
        title.text = meta.get("title", "Training Title")
        title.paragraphs[0].font.size = Pt(32)
        title.paragraphs[0].font.bold = True

        info_lines = [
            "SCHOOL GOVERNANCE AND OPERATIONS DIVISION (SGOD)",
            "School Management, Monitoring and Evaluation (SMME) Section",
            "QAME RESULT",
            "Republic of the Philippines · Department of Education · Region V · Schools Division of Masbate City",
            f"Title of Training: {meta.get('title','')}",
            f"Date Conducted: {meta.get('date','')}",
            f"Venue: {meta.get('venue','')}",
            f"Program Owner: {meta.get('owner','')}",
        ]
        box = slide.shapes.add_textbox(Inches(0.5), Inches(1.9), Inches(9), Inches(4.5)).text_frame
        for i, line in enumerate(info_lines):
            p = box.add_paragraph() if i else box.paragraphs[0]
            p.text = line
            p.font.size = Pt(14)
        return slide

    # ---- Category table slide
    def add_category_table_slide():
        if category_summary.empty:
            return
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        tf = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6)).text_frame
        tf.text = "EVALUATION RATING"
        tf.paragraphs[0].font.size = Pt(24)
        tf.paragraphs[0].font.bold = True

        # Order columns according to detected days
        file_cols = [c for c in category_summary.columns if c != "Overall Avg"]
        file_cols = _ordered_file_columns(category_summary, original_file_names)
        file_to_day = _map_file_to_day_label(original_file_names)

        headers = ["INDICATORS"] + [file_to_day.get(c, c) for c in file_cols] + ["Average"]

        rows = []
        for idx, row in category_summary.iterrows():
            r = [idx]
            for c in file_cols:
                r.append("" if pd.isna(row.get(c)) else f"{row.get(c):.2f}")
            r.append("" if pd.isna(row.get("Overall Avg")) else f"{row.get('Overall Avg'):.2f}")
            rows.append(r)

        n_rows = len(rows) + 1
        n_cols = len(headers)

        table = slide.shapes.add_table(n_rows, n_cols, Inches(0.5), Inches(1.2), Inches(9), Inches(4.5)).table
        for j, h in enumerate(headers):
            cell = table.cell(0, j)
            cell.text = h
            cell.text_frame.paragraphs[0].font.bold = True
        for i, data_row in enumerate(rows, start=1):
            for j, val in enumerate(data_row):
                table.cell(i, j).text = str(val)

        footer = slide.shapes.add_textbox(Inches(0.5), Inches(6.1), Inches(9), Inches(0.4)).text_frame
        footer.text = f"QAME Results · FY {pd.Timestamp.today().year}"

    # ---- Session tables (one slide per day)
    def add_session_slides():
        if session_summary.empty:
            return
        pattern = re.compile(r"DAY\s*(\d+)\s*-\s*LM\s*(\d+)", re.IGNORECASE)
        by_day: Dict[int, List[Tuple[str, float]]] = {}
        for s in session_summary.index:
            m = pattern.search(str(s))
            if not m:
                continue
            day_num, lm_num = int(m.group(1)), int(m.group(2))
            avg = session_summary.loc[s].get("Overall Avg", float("nan"))
            by_day.setdefault(day_num, []).append((f"DAY{day_num}-LM{lm_num}", avg))

        for day in sorted(by_day.keys()):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            title = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6)).text_frame
            title.text = "SESSION AVERAGES"
            title.paragraphs[0].font.size = Pt(24)
            title.paragraphs[0].font.bold = True

            sessions = sorted(by_day[day], key=lambda x: int(re.search(r"LM(\d+)", x[0]).group(1)))
            headers = ["SESSION TITLES", "RATING"]
            n_rows = len(sessions) + 1
            table = slide.shapes.add_table(n_rows, 2, Inches(0.5), Inches(1.2), Inches(9), Inches(4.5)).table
            table.cell(0, 0).text = headers[0]
            table.cell(0, 1).text = headers[1]
            for j in range(2):
                table.cell(0, j).text_frame.paragraphs[0].font.bold = True

            for i, (label, val) in enumerate(sessions, start=1):
                table.cell(i, 0).text = label
                table.cell(i, 1).text = "" if pd.isna(val) else f"{val:.2f}"

            footer = slide.shapes.add_textbox(Inches(0.5), Inches(6.1), Inches(9), Inches(0.4)).text_frame
            footer.text = f"QAME Results · FY {pd.Timestamp.today().year}"

    # ---- Overall slides (table + big number)
    def add_overall_slides():
        if not category_summary.empty:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            title = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6)).text_frame
            title.text = "OVERALL RATING"
            title.paragraphs[0].font.size = Pt(24)
            title.paragraphs[0].font.bold = True

            rows = []
            if "Overall Avg" in category_summary.columns:
                for cat, row in category_summary["Overall Avg"].items():
                    rows.append([cat, f"{row:.2f}"])
            avg_all_sessions = float("nan")
            if (not session_summary.empty) and ("Overall Avg" in session_summary.columns):
                avg_all_sessions = session_summary["Overall Avg"].mean()
                rows.append(["Average of All Sessions", f"{avg_all_sessions:.2f}"])

            overall_rating = float("nan")
            if "Overall Avg" in category_summary.columns:
                overall_rating = category_summary["Overall Avg"].mean()
            rows.append(["OVERALL RATING", f"{overall_rating:.2f}" if pd.notna(overall_rating) else ""])

            table = slide.shapes.add_table(len(rows) + 1, 2, Inches(0.5), Inches(1.2), Inches(6.5), Inches(4.5)).table
            table.cell(0, 0).text = "INDICATORS"
            table.cell(0, 1).text = "RATING"
            table.cell(0, 0).text_frame.paragraphs[0].font.bold = True
            table.cell(0, 1).text_frame.paragraphs[0].font.bold = True

            for i, (lbl, val) in enumerate(rows, start=1):
                table.cell(i, 0).text = str(lbl)
                table.cell(i, 1).text = str(val)

            footer = slide.shapes.add_textbox(Inches(0.5), Inches(6.1), Inches(9), Inches(0.4)).text_frame
            footer.text = f"QAME Results · FY {pd.Timestamp.today().year}"

        slide = prs.slides.add_slide(prs.slide_layouts[6])
        big = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(9), Inches(1.5)).text_frame
        if not category_summary.empty and "Overall Avg" in category_summary.columns:
            overall = category_summary["Overall Avg"].mean()
        else:
            overall = float("nan")
        big.text = f"{overall:.2f}" if pd.notna(overall) else "N/A"
        big.paragraphs[0].font.size = Pt(72)
        big.paragraphs[0].font.bold = True
        big.paragraphs[0].alignment = PP_ALIGN.CENTER

        sub = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(1)).text_frame
        sub.text = f"OUT OF {highest_score} AS THE HIGHEST SCORE"
        sub.paragraphs[0].font.size = Pt(20)
        sub.paragraphs[0].alignment = PP_ALIGN.CENTER

        quote = slide.shapes.add_textbox(Inches(0.5), Inches(4.6), Inches(9), Inches(1)).text_frame
        quote.text = "“OUTSTANDING”" if pd.notna(overall) and overall >= (0.9 * highest_score) else ""
        quote.paragraphs[0].font.size = Pt(24)
        quote.paragraphs[0].alignment = PP_ALIGN.CENTER
        quote.paragraphs[0].font.color.rgb = RGBColor(0, 128, 0)

    # Build slides
    add_title_slide()
    add_category_table_slide()
    add_session_slides()
    add_overall_slides()

    out = BytesIO()
    prs.save(out)
    out.seek(0)
    return out.getvalue()


def build_pdf_report(
    category_summary: pd.DataFrame,
    session_summary: pd.DataFrame,
    original_file_names: List[str],
    meta: Dict[str, str],
    highest_score: int = 5,
) -> bytes:
    """
    Simple PDF mirroring the PPT content (text + tables).
    """
    buf = BytesIO()
    doc = SimpleDocTemplate(buf, pagesize=A4, title=meta.get("title", "QAME Report"))
    styles = getSampleStyleSheet()
    flow = []

    # Title
    flow.append(Paragraph("<b>QAME RESULT</b>", styles["Title"]))
    flow.append(Paragraph(meta.get("title", ""), styles["Heading2"]))
    flow.append(Paragraph(f"Date: {meta.get('date','')} · Venue: {meta.get('venue','')} · Program Owner: {meta.get('owner','')}", styles["Normal"]))
    flow.append(Spacer(1, 12))

    # Categories table
    if not category_summary.empty:
        file_cols = [c for c in category_summary.columns if c != "Overall Avg"]
        file_cols = _ordered_file_columns(category_summary, original_file_names)
        file_to_day = _map_file_to_day_label(original_file_names)
        headers = ["INDICATORS"] + [file_to_day.get(c, c) for c in file_cols] + ["Average"]
        data = [headers]
        for idx, row in category_summary.iterrows():
            r = [idx] + [("" if pd.isna(row.get(c)) else f"{row.get(c):.2f}") for c in file_cols]
            r.append("" if pd.isna(row.get("Overall Avg")) else f"{row.get('Overall Avg'):.2f}")
            data.append(r)
        t = Table(data, hAlign="LEFT")
        t.setStyle(TableStyle([
            ("BACKGROUND", (0,0), (-1,0), colors.lightgrey),
            ("GRID", (0,0), (-1,-1), 0.5, colors.grey),
            ("FONTNAME", (0,0), (-1,0), "Helvetica-Bold"),
        ]))
        flow.append(Paragraph("<b>Evaluation Rating</b>", styles["Heading2"]))
        flow.append(t)
        flow.append(Spacer(1, 12))

    # Session tables per day
    if not session_summary.empty:
        pattern = re.compile(r"DAY\s*(\d+)\s*-\s*LM\s*(\d+)", re.IGNORECASE)
        by_day: Dict[int, List[Tuple[str, float]]] = {}
        for s in session_summary.index:
            m = pattern.search(str(s))
            if not m:
                continue
            dn, lm = int(m.group(1)), int(m.group(2))
            by_day.setdefault(dn, []).append((f"DAY{dn}-LM{lm}", session_summary.loc[s].get("Overall Avg", float("nan"))))
        for day in sorted(by_day.keys()):
            flow.append(Paragraph(f"<b>Session Averages — Day {day}</b>", styles["Heading2"]))
            data = [["SESSION TITLES", "RATING"]]
            for label, val in sorted(by_day[day], key=lambda x: int(re.search(r"LM(\d+)", x[0]).group(1))):
                data.append([label, "" if pd.isna(val) else f"{val:.2f}"])
            t = Table(data, hAlign="LEFT")
            t.setStyle(TableStyle([
                ("BACKGROUND",(0,0),(-1,0), colors.lightgrey),
                ("GRID",(0,0),(-1,-1), 0.5, colors.grey),
                ("FONTNAME",(0,0),(-1,0),"Helvetica-Bold")
            ]))
            flow.append(t)
            flow.append(Spacer(1, 8))

    # Overall
    overall = category_summary["Overall Avg"].mean() if (not category_summary.empty and "Overall Avg" in category_summary.columns) else float("nan")
    rows = []
    if not category_summary.empty and "Overall Avg" in category_summary.columns:
        for cat, val in category_summary["Overall Avg"].items():
            rows.append([cat, f"{val:.2f}"])
    avg_sessions = session_summary["Overall Avg"].mean() if (not session_summary.empty and "Overall Avg" in session_summary.columns) else float("nan")
    rows.append(["Average of All Sessions", "" if pd.isna(avg_sessions) else f"{avg_sessions:.2f}"])
    rows.append(["OVERALL RATING", "" if pd.isna(overall) else f"{overall:.2f}"])
    t = Table([["INDICATORS","RATING"]] + rows, hAlign="LEFT")
    t.setStyle(TableStyle([
        ("BACKGROUND",(0,0),(-1,0), colors.lightgrey),
        ("GRID",(0,0),(-1,-1), 0.5, colors.grey),
        ("FONTNAME",(0,0),(-1,0),"Helvetica-Bold")
    ]))
    flow.append(Paragraph("<b>Overall Rating</b>", styles["Heading2"]))
    flow.append(t)
    flow.append(Spacer(1, 12))
    flow.append(Paragraph(f"Out of {highest_score} as the highest score.", styles["Italic"]))

    doc.build(flow)
    buf.seek(0)
    return buf.getvalue()


# ──────────────────────────────────────────────────────────────────────────────
# UI: Upload (limit to 5 files)
# ──────────────────────────────────────────────────────────────────────────────
uploaded_files = st.file_uploader(
    "📂 Upload up to five CSV or Excel (.xlsx) files produced by the NEW TEMPLATE",
    type=["csv", "xlsx"],
    accept_multiple_files=True,
    help="You can mix CSV and XLSX. We'll read the first sheet if Excel has multiple sheets."
)

if uploaded_files:
    if len(uploaded_files) > 5:
        st.error(f"Please upload at most 5 files. You uploaded {len(uploaded_files)}.")
        st.stop()

    st.success(f"Loaded {len(uploaded_files)} file(s).")

    cat_tables: List[pd.DataFrame] = []
    ses_tables: List[pd.DataFrame] = []
    info_rows: List[Tuple[str, int]] = []

    for f in uploaded_files:
        try:
            df_raw = load_file(f)
        except Exception as e:
            st.error(f"❌ Could not read {f.name}: {e}")
            continue

        # Keep some quick info
        info_rows.append((f.name, len(df_raw)))

        # Compute per-file summaries
        cat_df = summarize_categories(df_raw, f.name)
        ses_df = summarize_sessions(df_raw, f.name)

        if not cat_df.empty:
            cat_tables.append(cat_df)
        if not ses_df.empty:
            ses_tables.append(ses_df)

    # ── Combine and display: Categories
    if cat_tables:
        st.subheader("📌 Category Averages (by file)")
        combined_cat = pd.concat(cat_tables, axis=1).sort_index()
        combined_cat = finalize_summary(combined_cat)
        combined_cat.index.name = "Category"

        st.dataframe(style_numeric(combined_cat), use_container_width=True)

        col1, col2 = st.columns(2)
        with col1:
            st.download_button(
                "⬇️ Download Category Summary (CSV)",
                data=to_csv_bytes(combined_cat),
                file_name="Category_Summary.csv",
                mime="text/csv",
            )
        with col2:
            st.download_button(
                "⬇️ Download Category Summary (Excel)",
                data=to_excel_bytes({"Category Summary": combined_cat}),
                file_name="Category_Summary.xlsx",
                mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            )

        # Chart
        melt_cat = combined_cat.reset_index().melt(
            id_vars="Category", var_name="File", value_name="Average"
        )
        # Exclude Overall Avg from the comparison bars
        melt_cat = melt_cat[melt_cat["File"] != "Overall Avg"]
        fig_cat = px.bar(
            melt_cat,
            x="Category",
            y="Average",
            color="File",
            barmode="group",
            title="Category Averages by File",
            text="Average",
        )
        fig_cat.update_traces(texttemplate="%{text:.2f}", textposition="outside")
        fig_cat.update_layout(yaxis_range=[0, 5])  # Likert 1–5 typical
        st.plotly_chart(fig_cat, use_container_width=True)
    else:
        st.info("No category columns found that match the NEW TEMPLATE prefixes.")

    # ── Combine and display: Sessions
    if ses_tables:
        st.subheader("📌 Session Averages (by file)")
        combined_ses = pd.concat(ses_tables, axis=1).sort_index()
        combined_ses = finalize_summary(combined_ses)
        combined_ses.index.name = "Session (DAY–LM)"

        st.dataframe(style_numeric(combined_ses), use_container_width=True)

        col1, col2 = st.columns(2)
        with col1:
            st.download_button(
                "⬇️ Download Session Summary (CSV)",
                data=to_csv_bytes(combined_ses),
                file_name="Session_Summary.csv",
                mime="text/csv",
            )
        with col2:
            st.download_button(
                "⬇️ Download Session Summary (Excel)",
                data=to_excel_bytes({"Session Summary": combined_ses}),
                file_name="Session_Summary.xlsx",
                mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            )

        # Chart
        melt_ses = combined_ses.reset_index().melt(
            id_vars="Session (DAY–LM)", var_name="File", value_name="Average"
        )
        melt_ses = melt_ses[melt_ses["File"] != "Overall Avg"]
        fig_ses = px.bar(
            melt_ses,
            x="Session (DAY–LM)",
            y="Average",
            color="File",
            barmode="group",
            title="Session Averages by File",
            text="Average",
        )
        fig_ses.update_traces(texttemplate="%{text:.2f}", textposition="outside")
        fig_ses.update_layout(yaxis_range=[0, 5])
        st.plotly_chart(fig_ses, use_container_width=True)
    else:
        st.info("No session columns (DAY–LM) detected in the uploaded files.")

    # Quick file info
    with st.expander("ℹ️ File statistics"):
        st.write(pd.DataFrame(info_rows, columns=["File", "Rows"]))

    # ──────────────────────────────────────────────────────────────────────────
    # NEW: Report Generation (PPTX & PDF) — mirrors your sample PPT format
    # ──────────────────────────────────────────────────────────────────────────
    st.markdown("### 📝 Generate Presentation Report")
    c1, c2 = st.columns(2)
    with c1:
        rpt_title = st.text_input("Title of Training", value="Revisiting and Rekindling Instructional Supervisory Practices (Under GABAY 2.0) Batch 3")
        rpt_date = st.text_input("Date Conducted", value="")
        rpt_venue = st.text_input("Venue", value="")
    with c2:
        rpt_owner = st.text_input("Program Owner", value="")
        highest_score = st.selectbox("Highest Score on the Scale", options=[5, 4], index=0)

    meta = {"title": rpt_title, "date": rpt_date, "venue": rpt_venue, "owner": rpt_owner}

    # Handle cases where tables may not exist (e.g., if none detected)
    categories_df = 'combined_cat' in locals() and isinstance(combined_cat, pd.DataFrame) and not combined_cat.empty
    sessions_df = 'combined_ses' in locals() and isinstance(combined_ses, pd.DataFrame) and not combined_ses.empty
    cat_df_for_report = combined_cat if categories_df else pd.DataFrame()
    ses_df_for_report = combined_ses if sessions_df else pd.DataFrame()

    file_names_in_order = [f.name for f in uploaded_files]  # day extraction will sort where applicable

    colA, colB = st.columns(2)
    with colA:
        ppt_bytes = build_pptx_report(cat_df_for_report, ses_df_for_report, file_names_in_order, meta, highest_score=highest_score)
        st.download_button(
            "📥 Download PowerPoint (PPTX)",
            data=ppt_bytes,
            file_name="QAME_Report.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            use_container_width=True,
        )
    with colB:
        pdf_bytes = build_pdf_report(cat_df_for_report, ses_df_for_report, file_names_in_order, meta, highest_score=highest_score)
        st.download_button(
            "📥 Download PDF",
            data=pdf_bytes,
            file_name="QAME_Report.pdf",
            mime="application/pdf",
            use_container_width=True,
        )

else:
    st.info("Upload up to five CSV/XLSX files to generate category and session summaries.")
